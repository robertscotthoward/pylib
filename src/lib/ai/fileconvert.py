"""
This module contains functions to convert documents to text, for example, converting a .doc file to a .docx file,
a .docx file to a text file, or a .pdf file to a text file.
"""

# uv add python-docx pypdf rdflib striprtf
import contextlib
import io
import re
import shutil
import tempfile
import threading
from docx import Document
from lib.tools import ensureFolder, readText, writeText
import ebooklib
import glob
import os
import subprocess
from lib.tools import *
import pandas as pd





def read_text_best_effort(filepath) -> str:
    """Read a text file whose encoding is unknown or already damaged.

    Tries UTF-8 first, then the Windows and Latin-1 code pages that produce mojibake
    rather than an error. clean_text() repairs whatever damage the fallback introduces.
    """
    with open(filepath, 'rb') as f:
        raw = f.read()
    for encoding in ('utf-8-sig', 'cp1252', 'latin-1'):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode('utf-8', errors='replace')


def is_nonsense_line(line: str, min_letters: int = 12) -> bool:
    """True when nostril judges a line to be gibberish, e.g. an OCR noise run.

    nostril is asked about the whole line rather than word by word. Its accuracy
    collapses at its six-character floor -- obvious junk like 'qwrtpz' and 'lkjhgf'
    each score as sensible in isolation, while the line containing them scores as
    nonsense -- so a line is the smallest unit worth asking about.

    Lines with fewer than `min_letters` letters are never judged. Short lines are
    where the detector is least reliable and where markdown syntax, headings, table
    rows, and code fragments live.
    """
    letters = re.sub(r"[^A-Za-z]", "", line)
    if len(letters) < min_letters:
        return False
    from nostril import nonsense
    try:
        return nonsense(line)
    except ValueError:
        # nostril declines to judge; keep the line rather than guess.
        return False


def clean_text(text: str, drop_nonsense: bool = True, min_letters: int = 12) -> str:
    """Clean up mangled text with ftfy and nostril.

    ftfy repairs the encoding layer: mojibake, doubly-encoded UTF-8, stray HTML
    entities, ligatures, curly-quote damage, and control characters. It also
    normalises line endings to \\n.

    nostril handles the content layer, dropping whole lines it judges to be gibberish,
    which is the shape OCR noise takes. Lines shorter than `min_letters` letters are
    left alone; see is_nonsense_line for why.

    Returns the cleaned text. Compare it against the input to tell whether anything changed.
    """
    import ftfy

    text = ftfy.fix_text(text)
    if not drop_nonsense:
        return text

    return '\n'.join(
        line for line in text.split('\n')
        if not is_nonsense_line(line, min_letters=min_letters)
    )


APPLE_DOUBLE_MAGIC = b'\x00\x05\x16\x07'


def is_apple_double(filepath) -> bool:
    """True for macOS AppleDouble sidecars, which carry a real document's name but not
    its content.

    Copying a folder from a Mac leaves a '._name.doc' beside every 'name.doc'. They hold
    resource-fork metadata, so every converter fails on them with an unhelpful error.
    The name is checked first because it costs nothing; the magic number confirms it.
    """
    if not os.path.basename(filepath).startswith('._'):
        return False
    try:
        with open(filepath, 'rb') as f:
            return f.read(4) == APPLE_DOUBLE_MAGIC
    except OSError:
        return False


def needs_conversion(source_path, target_path) -> bool:
    """True when target is missing or older than source, i.e. the target is stale.

    Equal timestamps count as up to date, so a target stamped with its source's
    mtime (see convert_doc_to_docx) is not rebuilt on every run.
    """
    if not os.path.exists(target_path):
        return True
    if not os.path.exists(source_path):
        return False
    return os.path.getmtime(source_path) > os.path.getmtime(target_path)


def get_text(filepath):
    """Get the text of a file. Only specific file extensions are supported."""

    try:
        if not os.path.exists(filepath):
            return None
        extension = os.path.splitext(filepath)[1].lower()
        if extension == ".pdf":
            return pdf_to_text(filepath)
        if extension == ".docx":
            return docx_to_text(filepath)
        if extension == ".rtf":
            return rtf_to_text(filepath)
        if extension == ".rdf":
            return rdf_to_text(filepath)
        if extension == ".epub":
            return epub_to_text(filepath)
        return readText(filepath)
    except Exception as e:
        print(f"Error getting text from {filepath}: {e}")
        return None




def get_markdown(filepath, repair_ocr: bool = True):
    """Get the text of a file. Only specific file extensions are supported.

    repair_ocr=False skips the LLM repair pass on OCR'd PDFs, letting a caller schedule
    those slow network calls itself. It has no effect on any other file type.
    """

    try:
        if not os.path.exists(filepath):
            return None
        # Matched on the lowercased extension: a folder copied from another machine is
        # full of .PDF and .DOC, and endswith('.pdf') silently sent those to readText.
        extension = os.path.splitext(filepath)[1].lower()
        if extension == ".pdf":
            return pdf_bytes_to_markdown(readBytes(filepath), repair_ocr=repair_ocr)
        if extension == ".docx":
            return docx_bytes_to_markdown(readBytes(filepath))
        if extension == ".rtf":
            return rtf_to_text(filepath)
        if extension == ".rdf":
            return rdf_to_text(filepath)
        if extension == ".epub":
            return epub_to_markdown(filepath)
        if extension == ".xlsx":
            return xlsx_bytes_to_markdown(readBytes(filepath))
        if extension == ".xls":
            return xls_bytes_to_markdown(readBytes(filepath))
        if extension == ".pptx":
            return pptx_bytes_to_markdown(readBytes(filepath))
        if extension == ".ppt":
            return ppt_bytes_to_markdown(readBytes(filepath))
        if extension in IMAGE_EXTENSIONS:
            return image_to_markdown(filepath)
        if extension in AUDIO_EXTENSIONS:
            return audio_to_markdown(filepath)
        return readText(filepath)
    except Exception as e:
        print(f"Error getting text from {filepath}: {e}")
        return None



def epub_to_markdown(filepath):
    import pypandoc
    return pypandoc.convert_file(filepath, 'md', format='epub')


def epub_to_text(filepath):
    from ebooklib import epub
    from bs4 import BeautifulSoup
    import re
    
    book = epub.read_epub(filepath)
    texts = []
    
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        # Get the content (HTML) from the EPUB item
        content = item.get_content()
        # Parse HTML and extract text
        soup = BeautifulSoup(content, 'html.parser')
        
        # Use space as separator to avoid breaking words
        text = soup.get_text(separator=' ', strip=True)
        
        # Clean up excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        if text:
            texts.append(text)
    
    text = '\n\n'.join(texts)
    return text




def rdf_to_text(filepath):
    import rdflib
    graph = rdflib.Graph()
    graph.parse(filepath, format='turtle')
    return graph.serialize(format='turtle') # Serialize to Turtle (very readable)




def rtf_to_text(filepath):
    from striprtf.striprtf import rtf_to_text
    s = readText(filepath)
    plain_text = rtf_to_text(s)
    return plain_text




def _find_pylib_file(filename):
    """Resolve a pylib config file (config.yaml / credentials.yaml).

    First searches up from the current working directory, so a consuming project's
    own config.yaml can override pylib's defaults (the documented convention). Falls
    back to pylib's own copy shipped alongside this module, so `pylib convert` still
    works when run against an arbitrary folder that isn't a pylib-consuming project
    (e.g. a folder of documents with no config.yaml anywhere in its ancestry).
    """
    path = findPath(filename, throwIfNotFound=False)
    if path:
        return path
    pylib_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
    fallback = os.path.join(pylib_root, filename)
    return fallback if os.path.exists(fallback) else None


def _ensure_tessdata_prefix():
    """Point PyMuPDF's built-in OCR at Tesseract's language data, so scanned/image-only
    PDFs (no text layer) get OCR'd instead of silently producing empty markdown.
    """
    if os.environ.get('TESSDATA_PREFIX'):
        return
    _locate_and_set_tessdata()


def _locate_and_set_tessdata():
    try:
        config_path = _find_pylib_file('config.yaml')
        config = getYaml(config_path) if config_path else {}
        tessdata_path = g(config, 'all/tessdata_path')
    except Exception:
        tessdata_path = None
    if not tessdata_path or not os.path.exists(tessdata_path):
        candidates = glob.glob(r"C:\Program Files\Tesseract-OCR\tessdata") + \
                     glob.glob(r"C:\Program Files (x86)\Tesseract-OCR\tessdata")
        tessdata_path = candidates[0] if candidates else None
    if tessdata_path and os.path.exists(tessdata_path):
        os.environ['TESSDATA_PREFIX'] = tessdata_path


AUDIO_EXTENSIONS = {'.mp3', '.wav', '.m4a', '.flac', '.ogg', '.oga', '.opus', '.aac', '.wma', '.aiff'}

# Whisper endpoints reject oversized uploads. Splitting long recordings needs ffmpeg, so
# for now an over-limit file is reported rather than silently half-transcribed.
AUDIO_MAX_MB = 25


def audio_to_markdown(filepath) -> str | None:
    """Transcribe an audio file to markdown with a speech-to-text model.

    Configured under config.yaml 'all/audio_transcribe'. Returns None when the feature is
    off, unconfigured, or the file is too large for the endpoint.
    """
    try:
        from lib.configurations import get_config_credentials_environment
        from lib.ai.modelstack import ModelStack

        config_path = _find_pylib_file('config.yaml')
        credentials_path = _find_pylib_file('credentials.yaml')
        config, _, _ = get_config_credentials_environment(config_path or 'config.yaml', credentials_path)
        settings = config.get('audio_transcribe') or {}
        if not settings.get('enabled', True) or not settings.get('model'):
            return None

        size_mb = os.path.getsize(filepath) / (1024 * 1024)
        limit = settings.get('max_mb', AUDIO_MAX_MB)
        if size_mb > limit:
            print(f"[WARN] {os.path.basename(filepath)} is {size_mb:.0f} MB, over the {limit} MB "
                  f"transcription limit; skipping. Split it or raise 'max_mb'.")
            return None

        model_config = settings['model']
        modelstack = ModelStack.from_config(model_config)
        with _llm_semaphore:
            text = modelstack.transcribe(
                readBytes(filepath),
                filename=os.path.basename(filepath),
                language=settings.get('language'),
            )
        if not text or not text.strip():
            return None
        return f"# {os.path.basename(filepath)}\n\n{text.strip()}\n"
    except Exception as e:
        print(f"[WARN] Transcription failed for {os.path.basename(filepath)}: {e}")
        return None


IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff', '.webp'}

_IMAGE_MIME = {
    '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
    '.gif': 'image/gif', '.bmp': 'image/bmp', '.tif': 'image/tiff',
    '.tiff': 'image/tiff', '.webp': 'image/webp',
}

# Longest edge, in pixels, an image is scaled to before upload. Vision models downsample
# anyway, and a modern camera JPEG inflates by a third as base64, so sending the original
# costs upload time and tokens without improving what the model can see.
IMAGE_MAX_EDGE = 1568


def _prepare_image(filepath) -> tuple:
    """Return (bytes, mime type) for upload, downscaling anything oversized to JPEG."""
    extension = os.path.splitext(filepath)[1].lower()
    mime = _IMAGE_MIME.get(extension, 'image/jpeg')
    raw = readBytes(filepath)

    try:
        from PIL import Image
        with Image.open(io.BytesIO(raw)) as img:
            if max(img.size) <= IMAGE_MAX_EDGE and extension in ('.jpg', '.jpeg', '.png'):
                return raw, mime
            img = img.convert('RGB')
            img.thumbnail((IMAGE_MAX_EDGE, IMAGE_MAX_EDGE))
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=85)
            return buffer.getvalue(), 'image/jpeg'
    except Exception as e:
        # An unreadable or exotic image still gets one attempt as-is.
        print(f"[WARN] Could not resize {os.path.basename(filepath)} ({e}); sending original.")
        return raw, mime


def image_to_markdown(filepath) -> str | None:
    """Describe an image with a vision model and return it as markdown.

    Configured under config.yaml 'all/image_describe'. Returns None when the feature is
    disabled or unconfigured, so callers can skip the file rather than record a failure.
    """
    try:
        from lib.configurations import get_config_credentials_environment
        from lib.ai.modelstack import ModelStack

        config_path = _find_pylib_file('config.yaml')
        credentials_path = _find_pylib_file('credentials.yaml')
        config, _, _ = get_config_credentials_environment(config_path or 'config.yaml', credentials_path)
        settings = config.get('image_describe') or {}
        if not settings.get('enabled', True) or not settings.get('model'):
            return None

        model_config = settings['model']
        prompt = settings.get('prompt', 'Describe this image in detail.')
        image_bytes, mime = _prepare_image(filepath)

        modelstack = ModelStack.from_config(model_config)
        with _llm_semaphore:
            description = modelstack.query_image(
                prompt, image_bytes, mime_type=mime,
                max_tokens=model_config.get('max_tokens', 2048),
            )
        if not description or not description.strip():
            return None
        return f"# {os.path.basename(filepath)}\n\n{description.strip()}\n"
    except Exception as e:
        print(f"[WARN] Image description failed for {os.path.basename(filepath)}: {e}")
        return None


def pdf_needs_ocr(filepath) -> bool:
    """True when a PDF has any page without a native text layer, i.e. OCR is required.

    Cheap: it reads the text objects already in the file and never runs OCR, so it is
    safe to call while deciding how to schedule work.
    """
    import pymupdf
    try:
        with pymupdf.open(filepath) as doc:
            if doc.page_count == 0:
                return False
            return any(not page.get_text().strip() for page in doc)
    except Exception:
        return False


# Total LLM repair calls allowed in flight at once, across every caller in the process.
# The repair endpoint slows down under load, and a slow response becomes a timeout, so
# concurrency has to be capped globally rather than per document: N documents each
# repairing M pages would otherwise put N*M requests on the wire.
LLM_CONCURRENCY = int(os.environ.get('PYLIB_LLM_CONCURRENCY', '4'))
_llm_semaphore = threading.BoundedSemaphore(LLM_CONCURRENCY)

# Above this many characters, one repair call risks exceeding max_tokens and truncating
# silently, so a long document is repaired page by page instead.
MAX_SINGLE_REPAIR_CHARS = 24000


def _repair_pages_concurrently(pages: dict) -> dict:
    """Repair several OCR'd pages at once, returning {index: repaired markdown}.

    Page repairs are independent network calls, so they run concurrently and are
    reassembled by index. The global semaphore still bounds what reaches the endpoint.
    """
    from concurrent.futures import ThreadPoolExecutor

    indexes = [i for i, md in pages.items() if md.strip()]
    if not indexes:
        return dict(pages)
    with ThreadPoolExecutor(max_workers=min(len(indexes), LLM_CONCURRENCY)) as pool:
        repaired = list(pool.map(reformat_ocr_markdown, [pages[i] for i in indexes]))
    result = dict(pages)
    result.update(dict(zip(indexes, repaired)))
    return result


def reformat_ocr_markdown(markdown: str) -> str:
    """Clean up Tesseract-OCR'd markdown by running it through a configured LLM.

    Configured via config.yaml 'all/ocr_postprocess' — swap models/providers there.
    'class: openai_compatible' works with any OpenAI-style /chat/completions API
    (DeepInfra, OpenRouter, Together, a local vLLM server, etc.), so switching
    providers is a config edit, not a code change. Falls back to the raw OCR
    markdown if the model call fails or post-processing is disabled.
    """
    try:
        from lib.configurations import get_config_credentials_environment
        from lib.ai.modelstack import ModelStack

        config_path = _find_pylib_file('config.yaml')
        credentials_path = _find_pylib_file('credentials.yaml')
        config, _, _ = get_config_credentials_environment(config_path or 'config.yaml', credentials_path)
        pp = config.get('ocr_postprocess') or {}
        if not pp.get('enabled', True):
            return markdown

        model_config = pp['model']
        prompt_prefix = pp.get('prompt', 'Format the following document to make it more readable:')
        modelstack = ModelStack.from_config(model_config)
        with _llm_semaphore:
            result = modelstack.query(f"{prompt_prefix}\n\n{markdown}",
                                      max_tokens=model_config.get('max_tokens', 8192))
        return result.strip() if result and result.strip() else markdown
    except Exception as e:
        print(f"[WARN] OCR post-processing formatting failed, using raw OCR markdown: {e}")
        return markdown


@contextlib.contextmanager
def _ocr_available(enabled: bool):
    """Turn PyMuPDF's built-in Tesseract OCR on or off for the duration of the block.

    MuPDF only OCRs when it can find Tesseract's language data, so TESSDATA_PREFIX is
    the on/off switch. Clearing it makes OCR impossible rather than merely unlikely,
    which is what keeps the fast extraction tiers fast: on a scanned page MuPDF returns
    empty in milliseconds instead of spending seconds in Tesseract.

    The previous value is always restored, so converting a folder of PDFs cannot leak
    an enabled OCR setting from one file into the next file's fast tier.
    """
    previous = os.environ.get('TESSDATA_PREFIX')
    try:
        if enabled:
            _ensure_tessdata_prefix()
        else:
            os.environ.pop('TESSDATA_PREFIX', None)
        yield
    finally:
        if previous is None:
            os.environ.pop('TESSDATA_PREFIX', None)
        else:
            os.environ['TESSDATA_PREFIX'] = previous


def _ocr_worker_count() -> int:
    """How many OCR processes to use. PYLIB_OCR_WORKERS overrides; 1 disables the pool."""
    override = os.environ.get('PYLIB_OCR_WORKERS')
    if override and override.isdigit() and int(override) > 0:
        return int(override)
    return min(8, os.cpu_count() or 1)


def _ocr_page_worker(task):
    """OCR a single page in a child process. Returns (page index, markdown).

    Module level and taking only picklable arguments, because ProcessPoolExecutor has to
    pickle it. Each child opens the PDF itself: PyMuPDF objects do not cross processes,
    and passing a path beats shipping the whole file to every worker.
    """
    pdf_path, index, tessdata = task
    if tessdata:
        os.environ['TESSDATA_PREFIX'] = tessdata
    import pymupdf
    import pymupdf4llm
    with pymupdf.open(pdf_path) as doc:
        chunks = pymupdf4llm.to_markdown(doc, pages=[index], page_chunks=True)
    return index, (chunks[0]['text'] if chunks else '')


def _ocr_pages(doc, pdf_bytes: bytes, pages, verbose: bool) -> dict:
    """OCR the given pages, in parallel processes when that is worth the startup cost.

    Tesseract is CPU-bound and single-threaded per page, so pages fan out across cores.
    Processes rather than threads because PyMuPDF is not thread-safe. One page, or a
    worker count of 1, runs inline instead: spawning a process costs about a second.
    """
    from concurrent.futures import ProcessPoolExecutor

    workers = min(_ocr_worker_count(), len(pages))
    if workers <= 1 or len(pages) <= 1:
        return _pdf_pages_to_markdown(doc, pages)

    tessdata = os.environ.get('TESSDATA_PREFIX')
    if verbose:
        print(f"  OCR: {len(pages)} pages across {workers} processes")
    try:
        with tempfile.TemporaryDirectory(prefix="pylib_ocr_") as staging:
            staged = os.path.join(staging, "doc.pdf")
            with open(staged, 'wb') as f:
                f.write(pdf_bytes)
            tasks = [(staged, index, tessdata) for index in sorted(pages)]
            with ProcessPoolExecutor(max_workers=workers) as pool:
                return dict(pool.map(_ocr_page_worker, tasks))
    except Exception as e:
        # A pool that cannot start is not a reason to lose the document.
        print(f"[WARN] Parallel OCR unavailable ({e}); falling back to sequential OCR.")
        return _pdf_pages_to_markdown(doc, pages)


def _pdf_pages_to_markdown(doc, pages) -> dict:
    """Convert a subset of pages, returned as {0-based page index: markdown}."""
    import pymupdf4llm

    if not pages:
        return {}
    chunks = pymupdf4llm.to_markdown(doc, pages=sorted(pages), page_chunks=True)
    return {chunk['metadata']['page_number'] - 1: chunk['text'] for chunk in chunks}


def pdf_bytes_to_markdown(bytes: bytes, allow_ocr: bool = True, verbose: bool = True,
                          repair_ocr: bool = True) -> str:
    """Convert a PDF to markdown, trying the fast extraction methods before OCR.

    Three tiers, cheapest first, and OCR is only ever reached for pages the cheap
    tiers could not read:

    1. PyMuPDF via pymupdf4llm on pages that carry a native text layer. Milliseconds
       per page, and it recovers headings, lists, and tables.
    2. pdfplumber over the whole document, used when a text layer exists but tier 1
       returned nothing -- a different parser sometimes reads what the layout analyser
       cannot.
    3. Tesseract OCR, restricted to the pages still empty after tiers 1 and 2, then
       passed through the configured LLM to repair OCR damage.

    Tiers 1 and 2 run with OCR switched off at the environment level, so a scanned page
    cannot silently pull Tesseract into the fast path. Pass allow_ocr=False to skip
    tier 3 entirely and accept whatever the fast tiers found.

    repair_ocr=False returns the raw OCR text without the LLM repair pass. OCR takes about
    a second while the repair is a network round trip an order of magnitude slower, so a
    caller converting many files can turn it off here and run the repairs concurrently
    (see pdf_needs_ocr and reformat_ocr_markdown).
    """
    import pymupdf

    doc = pymupdf.open(stream=io.BytesIO(bytes), filetype="pdf")
    total = doc.page_count
    native_pages = [i for i, page in enumerate(doc) if page.get_text().strip()]

    # Tier 1 -- native text layer only. OCR is switched off, so this cannot be slow.
    with _ocr_available(False):
        pages_md = _pdf_pages_to_markdown(doc, native_pages)

        missing = [i for i in range(total) if not pages_md.get(i, '').strip()]

        # Tier 2 -- a text layer exists but the layout analyser produced nothing from it.
        if missing and native_pages and not any(md.strip() for md in pages_md.values()):
            plumbed = pdf_bytes_to_text(bytes)
            if plumbed.strip():
                if verbose:
                    print(f"  PDF text: pdfplumber fallback ({total} pages)")
                return plumbed

    if not missing:
        if verbose:
            print(f"  PDF text: native text layer ({total} pages, no OCR needed)")
        return _join_pages(pages_md)

    if not allow_ocr:
        if verbose:
            print(f"  PDF text: native text layer, {len(missing)} of {total} pages skipped (OCR disabled)")
        return _join_pages(pages_md)

    # Tier 3 -- last resort, and only for the pages nothing else could read.
    if verbose:
        print(f"  PDF text: Tesseract OCR on {len(missing)} of {total} pages (slow path)")
    with _ocr_available(True):
        ocr_md = _ocr_pages(doc, bytes, missing, verbose)

    if not repair_ocr:
        pages_md.update(ocr_md)
        return _join_pages(pages_md)

    joined = _join_pages(ocr_md)
    if len(missing) == total and len(joined) <= MAX_SINGLE_REPAIR_CHARS:
        # Whole document scanned and short enough for one call: cheapest option, and it
        # gives the model the entire document as context.
        return reformat_ocr_markdown(joined) if joined.strip() else joined

    # Mixed document, or one long enough that a single call would risk truncation.
    # Repair page by page, concurrently, and reassemble in page order.
    pages_md.update(_repair_pages_concurrently(ocr_md))
    return _join_pages(pages_md)


def _join_pages(pages_md: dict) -> str:
    """Join per-page markdown in page order, dropping pages that came back empty."""
    return '\n\n'.join(
        pages_md[index] for index in sorted(pages_md) if pages_md[index].strip()
    )



def pdf_to_text(filepath):
    # import pypdf
    # pages = [page.extract_text() for page in pypdf.PdfReader(filepath).pages]
    # text = '\n\n'.join(pages)
    b = readBytes(filepath)
    return pdf_bytes_to_text(b)




def pdf_bytes_to_text(pdf_bytes: bytes) -> str:
    import pdfplumber
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        # extract_text returns None for a page with no text layer; '' keeps the join working.
        text = '\n\n'.join([page.extract_text(x_tolerance=5, y_tolerance=3, layout=True, x_density=7.25, y_density=13) or ''
                            for page in pdf.pages])
    return text




def convert_doc_to_docx(inPath, outPath=None):
    # Resolve Wordconv.exe: config.yaml > glob search
    wordconv_path = None
    try:
        config = getYaml('config.yaml') or {}
        wordconv_path = g(config, 'all/wordconv_path')
    except Exception:
        pass
    if not wordconv_path or not os.path.exists(wordconv_path):
        candidates = glob.glob(r"C:\Program Files\Microsoft Office*\**\Wordconv.exe", recursive=True)
        wordconv_path = candidates[0] if candidates else None
    if not wordconv_path or not os.path.exists(wordconv_path):
        print(f"FILE NOT FOUND: Wordconv.exe not found. Set 'all/wordconv_path' in config.yaml.")
        return None

    if not os.path.exists(inPath):
        print(f"Error: Source file not found at '{inPath}'")
        return None

    if is_apple_double(inPath):
        print(f"SKIPPED: '{os.path.basename(inPath)}' is a macOS resource fork, not a document.")
        return None

    if not outPath:
        outPath = os.path.splitext(inPath)[0] + ".docx"

    if os.path.exists(outPath):
        return outPath

    try:
        ensureFolder(os.path.dirname(outPath))
        print(f"Converting '{inPath}' to docx...")

        # Wordconv is a legacy MAX_PATH application: it fails with exit code -1 on any
        # path near or beyond 260 characters, even though Python and NTFS handle those
        # paths fine. Staging both sides through a short temp directory keeps Wordconv
        # well inside its limit no matter how deep the real file is buried.
        # ignore_cleanup_errors: Wordconv can outlive its own exit code and keep a handle
        # on the staging directory, which turns the teardown into a PermissionError.
        with tempfile.TemporaryDirectory(prefix="pylib_doc_", ignore_cleanup_errors=True) as staging:
            staged_in = os.path.join(staging, "in.doc")
            staged_out = os.path.join(staging, "out.docx")
            shutil.copyfile(inPath, staged_in)

            # No cwd: the staged paths are absolute, and making staging the child's
            # working directory is what stops it being deletable afterwards.
            subprocess.run([
                wordconv_path,
                "-oice",
                "-nme",
                staged_in,
                staged_out
            ], check=True, capture_output=True)

            if not os.path.exists(staged_out):
                print(f"CONVERSION FAILED: Wordconv reported success but wrote no file for '{inPath}'")
                return None

            shutil.move(staged_out, outPath)

        # Set the last updated time to the docx file
        os.utime(outPath, (os.path.getatime(inPath), os.path.getmtime(inPath)))
        return outPath

    except FileNotFoundError:
        print("CONVERSION FAILED: 'Wordconv.exe' not found. Set 'all/wordconv_path' in config.yaml.")
        return None
    except subprocess.CalledProcessError as e:
        detail = (e.stderr or b"").decode(errors="replace").strip() or f"exit code {e.returncode}"
        print(f"CONVERSION FAILED: Wordconv error on '{os.path.basename(inPath)}': {detail}")
        return None
    except OSError as e:
        print(f"CONVERSION FAILED: {type(e).__name__} on '{os.path.basename(inPath)}': {e}")
        return None




def convert_all_doc_to_docx(folder_path):
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.endswith(".doc"):
                convert_doc_to_docx(os.path.join(root, file))




def transform_all_doc_to_docx(inFolder, outFolder):
    for filepath in glob.glob(os.path.join(inFolder, '*.doc'), recursive=True):
        convert_doc_to_docx(filepath)




def docx_to_text(docx_path):
    from docx.table import Table

    document = Document(docx_path)
    parts = []
    for block in _iter_block_items(document):
        if isinstance(block, Table):
            lines = _table_to_markdown(block)
            if lines:
                parts.append('\n'.join(lines))
        elif block.text.strip():
            parts.append(block.text)
    return '\n\n'.join(parts)


def xls_bytes_to_markdown(byte_data):
    # 1. Wrap the byte array in a file-like object
    byte_stream = io.BytesIO(byte_data)
    
    # 2. Read the XLS file
    # Note: xlrd is required for .xls files
    df = pd.read_excel(byte_stream, engine='xlrd')
    
    # 3. Convert to Markdown
    # 'tabulate' is used under the hood for clean formatting
    return df.to_markdown(index=False)



def _iter_block_items(parent):
    """Yield Paragraph and Table objects from a document body or table cell, in document order.

    python-docx exposes `document.paragraphs` and `document.tables` as separate flat lists,
    which loses ordering and silently drops paragraphs nested inside table cells.
    """
    from docx.document import Document as _Document
    from docx.oxml.ns import qn
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    if isinstance(parent, _Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError(f"Cannot iterate block items of {type(parent)}")

    for child in parent_elm.iterchildren():
        if child.tag == qn('w:p'):
            yield Paragraph(child, parent)
        elif child.tag == qn('w:tbl'):
            yield Table(child, parent)


def _unique_row_cells(row):
    """Return a row's cells with horizontally/vertically merged duplicates collapsed."""
    try:
        cells = list(row.cells)
    except (IndexError, ValueError):
        # Irregular grids can break python-docx's cell mapping; fall back to raw <w:tc>.
        from docx.table import _Cell
        cells = [_Cell(tc, row.table) for tc in row._tr.tc_lst]

    unique = []
    seen = set()
    for cell in cells:
        key = id(cell._tc)
        if key in seen:
            continue
        seen.add(key)
        unique.append(cell)
    return unique


def _cell_content(cell, indent):
    """Return (inline_text, nested_table_lines) for a table cell."""
    from docx.table import Table

    texts = []
    nested = []
    for block in _iter_block_items(cell):
        if isinstance(block, Table):
            nested.extend(_table_to_markdown(block, indent + '  '))
        elif block.text.strip():
            texts.append(' '.join(block.text.split()))
    return ' '.join(texts), nested


def _table_to_markdown(table, indent=''):
    """Render a table as indented bullets: two-cell rows become `label: value`."""
    lines = []
    for row in table.rows:
        texts = []
        nested = []
        for cell in _unique_row_cells(row):
            text, cell_nested = _cell_content(cell, indent)
            texts.append(text)
            nested.extend(cell_nested)

        while texts and not texts[-1]:
            texts.pop()

        if len(texts) == 2:
            lines.append(f"{indent}- {texts[0]}: {texts[1]}")
        elif texts:
            lines.append(f"{indent}- " + ' | '.join(texts))
        lines.extend(nested)
    return lines


def _para_list_level(paragraph) -> int | None:
    """Return 0-based list indent level from XML numbering, or None if not a list paragraph."""
    from docx.oxml.ns import qn
    pPr = paragraph._p.find(qn('w:pPr'))
    if pPr is None:
        return None
    numPr = pPr.find(qn('w:numPr'))
    if numPr is None:
        return None
    ilvl = numPr.find(qn('w:ilvl'))
    if ilvl is None:
        return None
    try:
        return int(ilvl.get(qn('w:val'), 0))
    except (TypeError, ValueError):
        return None


def _paragraph_to_markdown(paragraph) -> str | None:
    """Render a single paragraph as markdown, or None when it is empty."""
    if not paragraph.text.strip():
        return None
    text = paragraph.text
    style = paragraph.style.name if paragraph.style else ''

    for n in range(1, 7):
        if style.startswith(f'Heading {n}'):
            return f"{'#' * n} {text}"

    level = _para_list_level(paragraph)
    if level is None and style.startswith('List'):
        level = 0
    if level is not None:
        return f"{'  ' * level}- {text}"
    return text


def docx_bytes_to_markdown(b : bytes) -> str | None:
    from docx.table import Table

    try:
        document = Document(io.BytesIO(b))
    except Exception as e:
        print(f"[ERROR] Failed to convert docx to markdown: {e}")
        return None
    markdown_lines = []

    for block in _iter_block_items(document):
        if isinstance(block, Table):
            lines = _table_to_markdown(block)
            if lines:
                markdown_lines.append('\n'.join(lines))
        else:
            md = _paragraph_to_markdown(block)
            if md is not None:
                markdown_lines.append(md)

    return '\n\n'.join(markdown_lines)




def doc_bytes_to_markdown(b : bytes) -> str:
    """
    Convert old Word (.doc) files to markdown text.
    Note: .doc files are binary format. This uses python-docx which has limited support.
    """
    return None
    inFile = r"cache/temp.doc"
    outFile = r"cache/temp.docx"
    writeBytes(inFile, b)
    convert_doc_to_docx(inFile, outFile)
    b =  readBytes(outFile)
    return docx_bytes_to_markdown(b)
    


def xlsx_bytes_to_markdown(b : bytes) -> str:
    import pandas as pd
    import io

    # We cannot just convert this to a markdown table because the columns are not always aligned and there migth be too many columns.
    excel_file = pd.ExcelFile(io.BytesIO(b))
    markdown_parts = []

    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(io.BytesIO(b), sheet_name=sheet_name)
        markdown_parts.append(f"# {sheet_name}")
        markdown_parts.append(df.to_csv(index=False, lineterminator='\n').strip('\n'))

    return '\n'.join(markdown_parts)




def pptx_bytes_to_markdown(b : bytes) -> str:
    import io
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(b))
    markdown_parts = []
    
    for slide_num, slide in enumerate(presentation.slides, 1):
        markdown_parts.append(f"## Slide {slide_num}\n")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                markdown_parts.append(shape.text)
    
    return '\n\n'.join(markdown_parts)




def ppt_bytes_to_markdown(b : bytes) -> str:
    """
    Convert old PowerPoint (.ppt) files to markdown text.
    Note: .ppt files are binary format and harder to parse than .pptx.
    This uses a basic extraction approach.
    """
    try:
        # Try using python-pptx with a workaround for older formats
        # If that fails, fall back to basic text extraction
        import io
        from pptx import Presentation
        
        try:
            presentation = Presentation(io.BytesIO(b))
            markdown_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                markdown_parts.append(f"## Slide {slide_num}\n")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        markdown_parts.append(shape.text)
            
            return '\n\n'.join(markdown_parts)
        except (KeyError, Exception):
            # If python-pptx fails, try basic binary text extraction
            text_content = []
            # Look for readable text in the binary data
            try:
                decoded = b.decode('utf-16-le', errors='ignore')
                # Extract text between common delimiters
                import re
                text_matches = re.findall(r'[\x20-\x7E]{4,}', decoded)
                text_content = [t.strip() for t in text_matches if t.strip() and len(t) > 3]
            except:
                pass
            
            if text_content:
                return '\n'.join(text_content)
            else:
                return "Could not extract text from .ppt file. The file may be corrupted or in an unsupported format."
    
    except Exception as e:
        print(f"Error converting PPT: {e}")
        return f"Error converting PPT file: {e}"


def all_files_to_text(folder_path, cleaned_extension='.cleaned', overwrite=False, filter=None):
    "For each file F in folder_path, clean F, convert F to text, and save the text to G where G = F + cleaned_extension"
    def keep(x): return True
    if filter is None:
        filter = keep
    for F in glob.glob(os.path.join(folder_path, '*.*'), recursive=True):
        if F.endswith(cleaned_extension):
            continue
        G = F + cleaned_extension
        # Rewrite G when it is missing or older than F.
        if overwrite or needs_conversion(F, G):
            text = get_text(F)
            # Replace all double spaces with single spaces
            lt = ""
            while lt != text:
                lt = text
                text = re.sub(r'\s+', ' ', text)
            if filter(text):
                writeText(G, text)
            else:
                if os.path.exists(G):
                    os.remove(G)
    pass




# ============================== TESTS ==============================

def test_convert_doc_to_docx():
    doc_file = r"..\data\corpus2\Niven, Larry - Unfinished Story.doc"
    extracted_text = convert_doc_to_docx(doc_file)
    print(extracted_text)

def test_convert_pdf_to_text():
    pdf_file = r"D:\rob\Wilmott Magazine\wilmott-202507-magazine-poulsen.pdf"
    extracted_text = get_text(pdf_file)
    print(extracted_text)


if __name__ == "__main__":
    # test_convert_doc_to_docx()
    test_convert_pdf_to_text()